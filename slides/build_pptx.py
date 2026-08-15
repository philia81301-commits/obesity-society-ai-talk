"""把離線 HTML 簡報轉成 PPTX（階段五）。

作法：Playwright 以 1280x720 視窗、2 倍解析度逐頁截圖，再用 python-pptx
把每張圖滿版貼進 16:9 投影片，並把該頁的文字放進「備忘稿」。

為什麼是截圖而不是重建原生元件：HTML 版才是母本，裡面有 SVG 圖標、
CSS 資料視覺化、手繪底圖與遮罩漸層，原生 PPTX 重畫不可能一致，
而這份 PPTX 的用途是「繳交存查與備援播放」。文字放進備忘稿，
所以內容仍可被搜尋與複製，不是純圖片黑盒。

用法：python build_pptx.py <輸出.pptx>
"""

import sys
import time
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

URL = "http://127.0.0.1:8765/slides/index.html"  # 需先啟動本機伺服器：python -m http.server 8765（在 repo 根目錄）
W, H = 1280, 720
SCALE = 2               # 2560x1440，投影與列印都夠銳利
SHOTS = Path("shots")


def render_slides() -> list[tuple[Path, str]]:
    """回傳 [(圖檔路徑, 該頁文字), ...]，依簡報順序。"""
    SHOTS.mkdir(exist_ok=True)
    out = []
    with sync_playwright() as pw:
        # 優先用系統 Chrome，省下另外下載 chromium
        try:
            browser = pw.chromium.launch(channel="chrome")
        except Exception:
            browser = pw.chromium.launch()
        page = browser.new_page(viewport={"width": W, "height": H},
                                device_scale_factor=SCALE)
        page.goto(URL, wait_until="networkidle")
        page.wait_for_function("typeof Reveal !== 'undefined' && Reveal.isReady()")

        # 匯出時隱藏操作介面；字型必須等載入完成，否則會截到 fallback 字型
        page.add_style_tag(content="""
            .reveal .controls, .reveal .progress, .reveal .slide-number { display:none !important; }
        """)
        page.evaluate("document.fonts.ready")
        time.sleep(1.2)

        total = page.evaluate("document.querySelectorAll('.slides > section').length")
        print(f"共 {total} 頁，開始截圖…")

        for i in range(total):
            page.evaluate(f"Reveal.slide({i})")
            time.sleep(0.45)          # 等 fade 轉場與背景切換
            num = page.evaluate(
                "document.querySelectorAll('.slides > section')[%d].getAttribute('data-slide')" % i
            )
            text = page.evaluate(
                "document.querySelectorAll('.slides > section')[%d].innerText" % i
            )
            f = SHOTS / f"slide_{i + 1:02d}.png"
            page.screenshot(path=str(f))
            out.append((f, (text or "").strip()))
            if (i + 1) % 10 == 0:
                print(f"  已完成 {i + 1}/{total}（P{num}）")

        browser.close()
    return out


def build(shots: list[tuple[Path, str]], dest: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)      # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]          # 空白版面

    for img, text in shots:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), Emu(0), Emu(0),
                             width=prs.slide_width, height=prs.slide_height)
        if text:
            # 備忘稿讓內容可被搜尋／複製，也方便大會存查時抽文字
            s.notes_slide.notes_text_frame.text = text

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))


def main() -> None:
    if len(sys.argv) < 2:
        sys.exit("用法：python build_pptx.py <輸出.pptx>")
    dest = Path(sys.argv[1])
    shots = render_slides()
    build(shots, dest)
    mb = dest.stat().st_size / 1024 / 1024
    print(f"\n已輸出：{dest}")
    print(f"頁數：{len(shots)}　檔案大小：{mb:.1f} MB")


if __name__ == "__main__":
    main()
